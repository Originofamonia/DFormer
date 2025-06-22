import os
from matplotlib import pyplot as plt
from matplotlib.colors import ListedColormap
import pickle
import matplotlib as mpl
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd
from pathlib import Path

import numpy as np


def compare_inferred_masks():
    model_paths = ['epoch-2_miou_85', 'epoch-3_miou_92', 'epoch-4_miou_94', 'epoch-5_miou_95']
    colors = ['#00000000', 'lime']
    cmap = ListedColormap(colors)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0.1)
    top_2 = Inches(6)
    width = Inches(14.0)
    height = Inches(1.2)
    alpha = 0.6
    df = pd.read_csv(f'datasets/trav/df2.csv', index_col=0)
    for row in df.itertuples(index=True, name='Pandas'):
        img_id = row.img.split('/')[-1].strip('.jpg')
        gt_path = row.img.replace('/images/', '/labels/')
        gt_file = os.path.splitext(gt_path)[0] + '.npy'
        ep_2_filename = f'output/{model_paths[0]}/{img_id}.npy'
        ep_3_filename = f'output/{model_paths[1]}/{img_id}.npy'
        ep_4_filename = f'output/{model_paths[2]}/{img_id}.npy'
        ep_5_filename = f'output/{model_paths[3]}/{img_id}.npy'
        slide = prs.slides.add_slide(blank_slide_layout)
        fig, axs = plt.subplots(2, 3, figsize=(14, 6))  # w,h
        # img, epoch2, epoch3
        #  gt, epoch4, epoch5
        q_img = plt.imread(row.img)
        q_target = np.load(gt_file)
        ep2 = np.load(ep_2_filename)
        ep3 = np.load(ep_3_filename)
        ep4 = np.load(ep_4_filename)
        ep5 = np.load(ep_5_filename)
        axs[0,0].imshow(q_img)
        axs[0,0].set_title(f'img')
        axs[0,0].axis('off')

        axs[1,0].imshow(q_img)
        axs[1,0].imshow(q_target, cmap=cmap, alpha=alpha)
        axs[1,0].set_title(f'target')
        axs[1,0].axis('off')

        axs[0,1].imshow(q_img)
        axs[0,1].imshow(ep2, cmap=cmap, alpha=alpha)
        axs[0,1].set_title(f'ep2')
        axs[0,1].axis('off')

        axs[0,2].imshow(q_img)
        axs[0,2].imshow(ep3, cmap=cmap, alpha=alpha)
        axs[0,2].set_title(f'ep3')
        axs[0,2].axis('off')

        axs[1,1].imshow(q_img)
        axs[1,1].imshow(ep4, cmap=cmap, alpha=alpha)
        axs[1,1].set_title(f'ep4')
        axs[1,1].axis('off')

        axs[1,2].imshow(q_img)
        axs[1,2].imshow(ep5, cmap=cmap, alpha=alpha)
        axs[1,2].set_title(f'ep5')
        axs[1,2].axis('off')

        plt.subplots_adjust(hspace=0.01, wspace=0.01)
        img_filename = f'output/pptx/{img_id}.png'
        fig.savefig(img_filename, bbox_inches='tight', pad_inches=0)
        plt.close()
        pic = slide.shapes.add_picture(img_filename, left, top)
        text_box = slide.shapes.add_textbox(left, top_2, width, height)

        # Get the text frame within the text box
        tf = text_box.text_frame

        # Add a paragraph to the text frame
        p = tf.add_paragraph()
        p.text = f'img: {row.img}'
    
    prs.save(f'output/pptx/different_models.pptx')


def draw_selected_0912():
    """
    TODO: random 1 row from df1 as query and random 1 row from df2 as support
    """
    model_paths = ['epoch-3_miou_92', 'epoch-5_miou_95']
    colors = ['#00000000', 'lime']
    cmap = ListedColormap(colors)
    dpi = 200
    sector_left = -45 #-135
    sector_right = 45 # 135
    angle_min = -26
    angle_max = 36
    angle_rad_min = np.deg2rad(angle_min)
    angle_rad_max = np.deg2rad(angle_max)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(0.1)
    top_2 = Inches(6)
    width = Inches(14.0)
    height = Inches(1.2)
    alpha = 0.5
    positives = [
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556012529532663.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556012529532663.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556009401066296.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661555916477569811.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556145691092960.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556022048261635.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661555874275943510.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556030035881989.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/positive/images/1661556172449825279.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661555934549742691.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/positive/images/1661555950491830103.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556184664382450.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556017621751062.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556310770125381.jpg',
    '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556188525106661.jpg',
    ]
    # negatives = [
    # '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556043648338549.jpg',
    # ]

    df1 = pd.read_csv(f'datasets/trav/df1.csv', index_col=0)  # s_imgs
    df2 = pd.read_csv(f'datasets/trav/df2.csv', index_col=0)

    for q_file in positives:
        img_id = q_file.split('/')[-1].strip('.jpg')
        q_gt_path = q_file.replace('/images/', '/labels/')
        q_gt_file = os.path.splitext(q_gt_path)[0] + '.npy'
        ep_3_filename = f'output/{model_paths[0]}/{img_id}.npy'
        ep_5_filename = f'output/{model_paths[1]}/{img_id}.npy'
        # q_files
        q_img = plt.imread(q_file)
        q_target = np.load(q_gt_file)
        q_laser_file = df2.loc[df2['img'] == q_file, 'laser'].values[0]
        with open(q_laser_file, 'rb') as q_f:
            data = pickle.load(q_f)
            q_laser = np.array(data['ranges'][::-1])[540:900]
        # s_files
        s_row = df1.sample(n=1)
        s_file, s_laser_file = s_row['img'].values[0], s_row['laser'].values[0]
        s_gt_path = s_file.replace('/images/', '/labels/')
        s_gt_file = os.path.splitext(s_gt_path)[0] + '.npy'
        with open(s_laser_file, 'rb') as s_f:
            data = pickle.load(s_f)
            s_laser = np.array(data['ranges'][::-1])[540:900]
        s_img = plt.imread(s_file)
        s_target = np.load(s_gt_file)
        ep3 = np.load(ep_3_filename)
        ep5 = np.load(ep_5_filename)
        slide = prs.slides.add_slide(blank_slide_layout)
        # s_img+s_tg, q_img+q_tg, q_before (ep3)
        # s_depth, q_depth, q_after (ep5)
        fig, axs = plt.subplots(2, 3, figsize=(14, 6))  # w,h
        axs[0,0].imshow(s_img)
        axs[0,0].imshow(s_target, cmap=cmap, alpha=alpha)
        axs[0,0].set_title(f's_img')
        axs[0,0].axis('off')

        angles = np.linspace(np.deg2rad(sector_left), np.deg2rad(sector_right), len(s_laser), endpoint=False)
        axs[1,0] = plt.subplot(234, projection='polar')
        axs[1,0].plot(angles, s_laser)
        axs[1,0].plot([angle_rad_max, angle_rad_max], [0, 5.1], color='red', linestyle='--')
        axs[1,0].plot([angle_rad_min, angle_rad_min], [0, 5.1], color='blue', linestyle='--')
        axs[1,0].set_thetamin(sector_left)
        axs[1,0].set_thetamax(sector_right)
        axs[1,0].set_theta_zero_location('N')
        axs[1,0].set_title(f"s_depth")
        axs[1,0].set_xticks(np.pi/180. * np.linspace(sector_left, sector_right, 10, endpoint=False))
        axs[1,0].figure.axes[3].set_axis_off()

        axs[0,1].imshow(q_img)
        axs[0,1].imshow(q_target, cmap=cmap, alpha=alpha)
        axs[0,1].set_title(f'q_img')
        axs[0,1].axis('off')

        angles = np.linspace(np.deg2rad(sector_left), np.deg2rad(sector_right), len(q_laser), endpoint=False)
        axs[1,1] = plt.subplot(235, projection='polar')
        axs[1,1].plot(angles, q_laser)
        axs[1,1].plot([angle_rad_max, angle_rad_max], [0, 5.1], color='red', linestyle='--')
        axs[1,1].plot([angle_rad_min, angle_rad_min], [0, 5.1], color='blue', linestyle='--')
        axs[1,1].set_thetamin(sector_left)
        axs[1,1].set_thetamax(sector_right)
        axs[1,1].set_theta_zero_location('N')
        axs[1,1].set_title(f"q_depth")
        axs[1,1].set_xticks(np.pi/180. * np.linspace(sector_left, sector_right, 10, endpoint=False))
        axs[1,1].figure.axes[4].set_axis_off()

        axs[0,2].imshow(q_img)
        axs[0,2].imshow(ep3, cmap=cmap, alpha=alpha)
        axs[0,2].set_title(f'q_before')
        axs[0,2].axis('off')

        axs[1,2].imshow(q_img)
        axs[1,2].imshow(ep5, cmap=cmap, alpha=alpha)
        axs[1,2].set_title(f'q_after')
        axs[1,2].axis('off')

        plt.subplots_adjust(hspace=0.15, wspace=0.01)
        img_filename = f'output/0912/{img_id}.png'
        fig.savefig(img_filename,bbox_inches='tight', pad_inches=0.01, dpi=dpi)
        plt.close()
        pic = slide.shapes.add_picture(img_filename, left, top)

    prs.save(f'output/pptx/0912.pptx')


def draw_arch_sucai():
    """
    Need q_img, q_depth, q_mask/q_pred
    """
    # colors = ['#00000000', 'lime']
    colors = ['darkgray', 'lime']
    cmap = ListedColormap(colors)
    alpha = 0.8
    dpi = 200
    sector_left = -45 #-135
    sector_right = 45 # 135
    angle_min = -26
    angle_max = 36
    angle_rad_min = np.deg2rad(angle_min)
    angle_rad_max = np.deg2rad(angle_max)
    df2 = pd.read_csv(f'datasets/trav/df2.csv', index_col=0)
    q_file = '/home/qiyuan/2023spring/segmentation_indoor_images/uc/challenging/images/1661556184664382450.jpg'
    img_id = q_file.split('/')[-1].strip('.jpg')
    q_gt_path = q_file.replace('/images/', '/labels/')
    q_gt_file = os.path.splitext(q_gt_path)[0] + '.npy'
    # q_img = plt.imread(q_file)
    q_target = np.load(q_gt_file)
    # q_laser_file = df2.loc[df2['img'] == q_file, 'laser'].values[0]
    # with open(q_laser_file, 'rb') as q_f:
    #     data = pickle.load(q_f)
    #     q_laser = np.array(data['ranges'][::-1])[540:900]
    # angles = np.linspace(np.deg2rad(sector_left), np.deg2rad(sector_right), len(q_laser), endpoint=False)
    # fig, axs = plt.subplots(1, 1, figsize=(4, 3))
    # axs[0,0] = plt.subplot(111, projection='polar')
    # axs[0,0].plot(angles, q_laser)
    # axs[0,0].plot([angle_rad_max, angle_rad_max], [0, 5.1], color='red', linestyle='--')
    # axs[0,0].plot([angle_rad_min, angle_rad_min], [0, 5.1], color='blue', linestyle='--')
    # axs[0,0].set_thetamin(sector_left)
    # axs[0,0].set_thetamax(sector_right)
    # axs[0,0].set_theta_zero_location('N')
    # # axs[0,0].set_title(f"s_depth")
    # axs[0,0].set_xticks(np.pi/180. * np.linspace(sector_left, sector_right, 10, endpoint=False))
    # axs[0,0].figure.axes[0].set_axis_off()

    # Plot the image
    # plt.imshow(array_2d, cmap='gray', interpolation='none')
    plt.imshow(q_target, cmap=cmap, alpha=alpha)
    # Remove axis and borders
    plt.axis('off')  # Hides the axis
    plt.gca().set_position([0, 0, 1, 1])  # Removes padding around the image
    img_filename = f'output/0912/{img_id}_mask.png'
    plt.savefig(img_filename,bbox_inches='tight',pad_inches=0.0, dpi=dpi)
    plt.close()
    


if __name__ == '__main__':
    # compare_inferred_masks()
    # draw_selected_0912()
    draw_arch_sucai()
